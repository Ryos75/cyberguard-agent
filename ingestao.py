from dotenv import load_dotenv
load_dotenv()

import json
from pathlib import Path
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, CSVLoader, BSHTMLLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_chroma import Chroma

def load_xlsx(path):
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True); docs = []
    for ws in wb.worksheets:
        linhas = [[str(c.value) if c.value is not None else "" for c in row] for row in ws.iter_rows()]
        linhas = [l for l in linhas if any(l)]
        if not linhas: continue
        header = linhas[0]
        for row in linhas[1:]:
            pares = " | ".join(f"{h}: {v}" for h, v in zip(header, row) if v)
            if pares.strip():
                docs.append(Document(page_content=pares, metadata={"fonte": Path(path).name}))
    return docs

def load_pptx(path):
    from pptx import Presentation
    prs = Presentation(path); docs = []
    for i, slide in enumerate(prs.slides, 1):
        txt = "\n".join(sh.text for sh in slide.shapes if sh.has_text_frame and sh.text.strip())
        if txt.strip():
            docs.append(Document(page_content=f"[Slide {i}]\n{txt}", metadata={"fonte": Path(path).name}))
    return docs

def load_json(path):
    with open(path, encoding="utf-8") as f: dados = json.load(f)
    return [Document(page_content=f"{i['termo']}: {i['definicao']}", metadata={"fonte": Path(path).name}) for i in dados.get("glossario", [])]

def load_md(path):
    return [Document(page_content=Path(path).read_text(encoding="utf-8"), metadata={"fonte": Path(path).name})]

LOADERS = {
    ".pdf": lambda p: PyPDFLoader(p).load(), ".docx": lambda p: Docx2txtLoader(p).load(),
    ".csv": lambda p: CSVLoader(p).load(), ".html": lambda p: BSHTMLLoader(p).load(),
    ".xlsx": load_xlsx, ".pptx": load_pptx, ".json": load_json, ".md": load_md,
}

def carregar(pasta="docs"):
    docs = []
    for arq in sorted(Path(pasta).iterdir()):
        ext = arq.suffix.lower()
        if ext not in LOADERS: continue
        try:
            d = LOADERS[ext](str(arq))
            for x in d: x.metadata.setdefault("fonte", arq.name)
            docs.extend(d); print(f"OK {arq.name}: {len(d)} doc(s)")
        except Exception as e: print(f"ERRO {arq.name}: {e}")
    return docs

if __name__ == "__main__":
    docs = carregar()
    chunks = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200).split_documents(docs)
    print(f"{len(chunks)} chunks gerados")
    Chroma.from_documents(chunks, GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001"), persist_directory="./chroma_db")
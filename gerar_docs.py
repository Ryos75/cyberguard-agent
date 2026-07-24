"""Gera os documentos fictícios de Segurança da Informação da TechCorp."""
import os, json, csv

os.makedirs("docs", exist_ok=True)

# ---------- 1. PDF (politica_seguranca.pdf) ----------
from fpdf import FPDF

texto_pdf = """POLITICA DE SEGURANCA DA INFORMACAO (PSI) - TECHCORP
Versao 5.1 | Vigencia: 01/01/2025 | Classificacao: Interna

1. OBJETIVO
Estabelecer diretrizes de seguranca da informacao para todos os
colaboradores, prestadores e parceiros da TechCorp Tecnologia Ltda.

2. CLASSIFICACAO DA INFORMACAO
- PUBLICA: pode ser divulgada livremente (ex: material de marketing).
- INTERNA: uso interno, sem dano grave se vazada (ex: comunicados).
- CONFIDENCIAL: dano significativo se vazada (ex: contratos, dados
  de clientes). Requer criptografia em transito e em repouso.
- RESTRITA: dano severo (ex: credenciais, dados financeiros, codigo-fonte
  de produtos). Acesso somente com aprovacao do gestor + CISO.

3. AUTENTICACAO
- MFA (autenticacao multifator) e OBRIGATORIO em todos os sistemas
  corporativos, sem excecoes.
- Detalhes de senhas: consultar o Manual de Senhas e MFA.

4. CONTROLE DE ACESSO
- Principio do menor privilegio: cada colaborador acessa apenas o
  necessario para sua funcao.
- Revisao de acessos: trimestral, conduzida pelos gestores.
- Desligamento: acessos revogados em ate 4 horas apos o desligamento.

5. EQUIPE DE SEGURANCA (SOC)
- E-mail: soc@techcorp.com | Ramal: 4500 (24x7)
- CISO: Ricardo Tavares | security@techcorp.com
- Reporte de incidentes: IMEDIATO, em ate 15 minutos apos a suspeita.

6. USO ACEITAVEL
- Proibido instalar software nao homologado (shadow IT).
- Proibido usar e-mail corporativo em servicos pessoais.
- Dispositivos USB: somente os criptografados e fornecidos pela empresa.

7. BACKUPS
- Dados corporativos: backup diario automatico, retencao de 90 dias.
- Testes de restauracao: mensais, responsabilidade da equipe de TI.

8. SANCOES
Descumprimento sujeita o colaborador a advertencia, suspensao ou
desligamento por justa causa, conforme gravidade.

9. TREINAMENTO
Treinamento de conscientizacao: obrigatorio no onboarding e reciclagem
anual. Simulacoes de phishing: trimestrais."""

pdf = FPDF()
pdf.add_page()
pdf.set_font("Helvetica", size=11)
pdf.multi_cell(0, 6, texto_pdf)
pdf.output("docs/politica_seguranca.pdf")
print("✅ politica_seguranca.pdf")

# ---------- 2. DOCX (manual_senhas_mfa.docx) ----------
from docx import Document

doc = Document()
doc.add_heading("MANUAL DE SENHAS E MFA — TECHCORP", 0)
doc.add_paragraph("Documento: SEC-MAN-002 | Revisão: 3 | Aprovado por: CISO Ricardo Tavares")
secoes = [
    ("1. REGRAS DE SENHA",
     "Mínimo de 12 caracteres, combinando letras maiúsculas, minúsculas, "
     "números e símbolos. Troca obrigatória a cada 180 dias. É proibido "
     "reutilizar qualquer uma das últimas 5 senhas. Senhas não podem conter "
     "nome, data de nascimento ou o nome da empresa."),
    ("2. GERENCIADOR DE SENHAS",
     "O gerenciador oficial e homologado da TechCorp é o Bitwarden (licença "
     "corporativa). É PROIBIDO salvar senhas em planilhas, blocos de notas, "
     "post-its ou no navegador."),
    ("3. MFA — AUTENTICAÇÃO MULTIFATOR",
     "Obrigatório em todos os sistemas. Método padrão: aplicativo Microsoft "
     "Authenticator. SMS só é permitido como método de backup. Nunca aprove "
     "uma notificação de MFA que você não solicitou — isso é sinal de ataque "
     "(MFA fatigue). Se ocorrer, troque a senha imediatamente e avise o SOC "
     "no ramal 4500."),
    ("4. COMPARTILHAMENTO",
     "Senhas são pessoais e intransferíveis. NINGUÉM da TI ou da diretoria "
     "jamais pedirá sua senha, por nenhum canal. Compartilhamento de senha "
     "é violação grave da PSI."),
    ("5. SENHAS COMPROMETIDAS",
     "Se suspeitar que sua senha vazou: 1) troque imediatamente; "
     "2) avise o SOC (soc@techcorp.com / ramal 4500); 3) verifique "
     "atividades recentes na conta. Prazo para reporte: 15 minutos."),
]
for titulo, texto in secoes:
    doc.add_heading(titulo, level=2)
    doc.add_paragraph(texto)
doc.save("docs/manual_senhas_mfa.docx")
print("✅ manual_senhas_mfa.docx")

# ---------- 3. XLSX (inventario_ativos.xlsx) ----------
from openpyxl import Workbook

wb = Workbook()
ws = wb.active
ws.title = "Ativos Críticos"
ws.append(["ID", "Ativo", "Tipo", "Criticidade", "Dono (Owner)",
           "RTO", "RPO", "Backup", "Exposto à Internet"])
ativos = [
    ["AT-001", "ERP Financeiro", "Sistema", "Crítica", "Carlos Mendes", "4 horas", "1 hora", "Diário", "Não"],
    ["AT-002", "VPN Corporativa", "Infraestrutura", "Crítica", "Equipe TI", "2 horas", "N/A", "Config semanal", "Sim"],
    ["AT-003", "E-mail (M365)", "Sistema", "Crítica", "Equipe TI", "6 horas", "30 min", "Nativo M365", "Sim"],
    ["AT-004", "CRM Vendas", "Sistema", "Alta", "Fernanda Rocha", "8 horas", "4 horas", "Diário", "Sim"],
    ["AT-005", "Servidor de Arquivos", "Infraestrutura", "Alta", "Equipe TI", "12 horas", "12 horas", "Diário", "Não"],
    ["AT-006", "Site Institucional", "Sistema", "Média", "Marketing", "24 horas", "24 horas", "Semanal", "Sim"],
    ["AT-007", "Repositório de Código (GitLab)", "Sistema", "Crítica", "Ana Paula Silva", "4 horas", "1 hora", "Diário", "Não"],
    ["AT-008", "Wi-Fi Corporativo", "Infraestrutura", "Média", "Equipe TI", "8 horas", "N/A", "Config semanal", "Não"],
]
for linha in ativos:
    ws.append(linha)
wb.save("docs/inventario_ativos.xlsx")
print("✅ inventario_ativos.xlsx")

# ---------- 4. PPTX (treinamento_phishing.pptx) ----------
from pptx import Presentation

prs = Presentation()
slides = [
    ("Treinamento Anti-Phishing TechCorp 2025",
     "Obrigatório para todos os colaboradores\nReciclagem anual + simulações trimestrais"),
    ("O que é Phishing?",
     "Golpe que usa e-mails, SMS ou ligações falsas para roubar credenciais,"
     "\ndados ou instalar malware.\n"
     "91% dos ciberataques começam com um e-mail de phishing."),
    ("Os 7 Sinais de um E-mail de Phishing",
     "1. Senso de urgência exagerado ('sua conta será bloqueada HOJE')\n"
     "2. Remetente estranho ou domínio parecido (techc0rp.com)\n"
     "3. Erros de português ou formatação\n"
     "4. Links que não batem com o texto (passe o mouse antes de clicar!)\n"
     "5. Anexos inesperados (.zip, .exe, .html)\n"
     "6. Pedido de dados sensíveis ou senha\n"
     "7. Ofertas boas demais para ser verdade"),
    ("Recebeu um e-mail suspeito?",
     "NÃO clique em links nem abra anexos.\n"
     "Use o botão 'Reportar Phishing' no Outlook.\n"
     "Na dúvida, encaminhe para soc@techcorp.com ou ligue no ramal 4500."),
    ("Clicou sem querer?",
     "Acontece! O importante é agir rápido:\n"
     "1. NÃO desligue o computador\n"
     "2. Desconecte o cabo de rede / Wi-Fi\n"
     "3. Avise o SOC IMEDIATAMENTE (ramal 4500) — prazo: 15 minutos\n"
     "4. Não tente resolver sozinho\n"
     "Reportar rápido NUNCA gera punição. Esconder o incidente, sim."),
    ("Simulações Trimestrais",
     "A TechCorp envia e-mails de phishing SIMULADOS a cada trimestre.\n"
     "Quem clicar recebe orientação na hora (sem punição).\n"
     "Quem clicar 2 vezes no mesmo ano refaz o treinamento completo."),
    ("Certificação",
     "Quiz na plataforma UniTechCorp.\nNota mínima: 85%.\nCertificado válido por 12 meses."),
]
layout = prs.slide_layouts[1]
for titulo, corpo in slides:
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = titulo
    slide.placeholders[1].text = corpo
prs.save("docs/treinamento_phishing.pptx")
print("✅ treinamento_phishing.pptx")

# ---------- 5. MD (politica_home_office.md) ----------
politica_md = """# Política de Segurança no Trabalho Remoto — TechCorp
> Versão 2.4 — Aprovada pelo CISO em 10/01/2025

## 1. VPN Obrigatória
- Todo acesso a sistemas corporativos fora do escritório DEVE usar a VPN.
- A VPN desconecta após 12 horas; reconecte ao iniciar o expediente.
- Problemas com VPN: chamado no helpdesk ou ramal 4400.

## 2. Redes Wi-Fi
- **Proibido** usar Wi-Fi público (cafés, aeroportos, hotéis) sem VPN ativa.
- Rede doméstica: troque a senha padrão do roteador e use WPA2 ou WPA3.
- Nunca compartilhe a rede do trabalho com visitas em casa.

## 3. Equipamentos
- Use SOMENTE o notebook corporativo para trabalho.
- **Proibido** conectar dispositivos pessoais (pendrive, HD externo) ao
  equipamento corporativo.
- Bloqueie a tela SEMPRE que se afastar: tecle **Win + L** (leva 1 segundo!).
- Familiares não devem usar o equipamento corporativo, nem rapidamente.

## 4. Ambiente Físico
- Evite trabalhar com dados confidenciais em locais públicos.
- Use película de privacidade em viagens (solicite ao TI).
- Reuniões sensíveis: use fone de ouvido.

## 5. Videoconferências
- Não compartilhe links de reunião em redes sociais.
- Ative a sala de espera para reuniões com externos.
- Cuidado com o que aparece ao compartilhar a tela (feche e-mails e chats).

## 6. Perda ou Roubo de Equipamento
- Comunique o SOC IMEDIATAMENTE: soc@techcorp.com ou ramal 4500 (24x7).
- O bloqueio e a limpeza remota do equipamento ocorrem em até 2 horas.
- Registre boletim de ocorrência e envie ao RH em até 48 horas.
"""
with open("docs/politica_home_office.md", "w", encoding="utf-8") as f:
    f.write(politica_md)
print("✅ politica_home_office.md")

# ---------- 6. CSV (historico_incidentes.csv) ----------
incidentes = [
    ["id", "data", "tipo", "severidade", "descricao", "causa_raiz", "licao_aprendida", "status"],
    ["INC-2024-012", "05/02/2024", "Malware", "Média", "Notebook com adware após download de software não homologado", "Shadow IT", "Reforço da política de software homologado", "Encerrado"],
    ["INC-2024-031", "18/04/2024", "Phishing", "Alta", "Colaborador inseriu credenciais em página falsa do M365; conta comprometida enviou 200 e-mails", "E-mail de phishing com domínio parecido", "MFA bloqueou acesso externo; campanha de conscientização extra", "Encerrado"],
    ["INC-2024-045", "22/06/2024", "Mídia removível", "Média", "Pendrive pessoal infectado conectado a notebook corporativo", "Violação da política de USB", "Bloqueio técnico de portas USB para mídias não criptografadas", "Encerrado"],
    ["INC-2024-058", "10/08/2024", "Perda de equipamento", "Alta", "Notebook esquecido em táxi durante viagem", "Descuido físico", "Limpeza remota executada em 40 min; sem vazamento confirmado", "Encerrado"],
    ["INC-2024-071", "03/10/2024", "MFA Fatigue", "Alta", "Colaborador aprovou notificação MFA não solicitada de madrugada", "Ataque de MFA fatigue após vazamento de senha", "Migração de SMS para app authenticator; alerta sobre aprovações indevidas", "Encerrado"],
    ["INC-2024-089", "27/11/2024", "Engenharia social", "Média", "Ligação se passando pela TI pedindo senha; colaborador reportou sem fornecer", "Vishing", "Caso usado como exemplo positivo no treinamento", "Encerrado"],
    ["INC-2025-003", "15/01/2025", "Phishing", "Baixa", "E-mail de phishing reportado por 14 colaboradores em 10 minutos", "Campanha externa em massa", "Tempo de reporte recorde; bloqueio do remetente em 8 min", "Encerrado"],
]
with open("docs/historico_incidentes.csv", "w", newline="", encoding="utf-8") as f:
    csv.writer(f).writerows(incidentes)
print("✅ historico_incidentes.csv")

# ---------- 7. JSON (glossario_seguranca.json) ----------
glossario = {"glossario": [
    {"termo": "Phishing", "definicao": "Golpe por e-mail, SMS ou ligação que imita fontes confiáveis para roubar credenciais ou instalar malware. Na TechCorp, reporte pelo botão 'Reportar Phishing' no Outlook ou ao SOC (ramal 4500)."},
    {"termo": "Ransomware", "definicao": "Malware que criptografa arquivos e exige resgate. A TechCorp mantém backups diários com retenção de 90 dias como proteção."},
    {"termo": "MFA", "definicao": "Autenticação multifator: exige um segundo fator além da senha. Obrigatório em todos os sistemas da TechCorp. Método padrão: Microsoft Authenticator."},
    {"termo": "MFA Fatigue", "definicao": "Ataque em que o criminoso dispara várias notificações de MFA esperando que a vítima aprove por cansaço. Nunca aprove notificações que você não solicitou."},
    {"termo": "VPN", "definicao": "Rede privada virtual que criptografa a conexão. Obrigatória para qualquer acesso remoto a sistemas da TechCorp."},
    {"termo": "Engenharia Social", "definicao": "Manipulação psicológica para obter informações ou acessos. Inclui phishing (e-mail), vishing (ligação) e smishing (SMS)."},
    {"termo": "Shadow IT", "definicao": "Uso de softwares ou serviços não homologados pela TI. Proibido na TechCorp por criar riscos invisíveis de segurança."},
    {"termo": "Zero Trust", "definicao": "Modelo de segurança que não confia em nada por padrão: toda conexão é verificada, independentemente de estar dentro da rede."},
    {"termo": "BYOD", "definicao": "Bring Your Own Device (uso de dispositivo pessoal). Na TechCorp, dispositivos pessoais NÃO podem acessar a rede corporativa nem ser conectados a equipamentos da empresa."},
    {"termo": "SOC", "definicao": "Security Operations Center: equipe que monitora e responde a incidentes 24x7. Contato: soc@techcorp.com ou ramal 4500."},
    {"termo": "RTO", "definicao": "Recovery Time Objective: tempo máximo aceitável para restaurar um sistema após falha. Ex.: a VPN da TechCorp tem RTO de 2 horas."},
    {"termo": "Vishing", "definicao": "Phishing por voz (ligação telefônica). Lembre-se: NINGUÉM da TI da TechCorp pede senha por telefone, jamais."},
]}
with open("docs/glossario_seguranca.json", "w", encoding="utf-8") as f:
    json.dump(glossario, f, ensure_ascii=False, indent=2)
print("✅ glossario_seguranca.json")

# ---------- 8. HTML (faq_incidentes.html) ----------
faq_html = """<!DOCTYPE html>
<html lang="pt-BR">
<head><meta charset="UTF-8"><title>FAQ — O que fazer se... | Segurança TechCorp</title></head>
<body>
<h1>FAQ de Resposta a Incidentes — "O que fazer se..."</h1>

<h2>1. Cliquei em um link suspeito. E agora?</h2>
<p>NÃO desligue o computador (isso destrói evidências). Desconecte o cabo de
rede ou o Wi-Fi e avise o SOC IMEDIATAMENTE: soc@techcorp.com ou ramal 4500.
Prazo para reporte: 15 minutos. Reportar rápido nunca gera punição.</p>

<h2>2. Digitei minha senha em um site falso.</h2>
<p>Troque a senha imediatamente em todos os sistemas onde ela é usada,
avise o SOC (ramal 4500) e fique atento a notificações de MFA não
solicitadas — não as aprove.</p>

<h2>3. Perdi meu notebook / fui roubado.</h2>
<p>Comunique o SOC imediatamente (24x7): soc@techcorp.com ou ramal 4500.
O bloqueio e a limpeza remota ocorrem em até 2 horas. Depois, registre
boletim de ocorrência e envie ao RH em até 48 horas.</p>

<h2>4. Recebi uma ligação da "TI" pedindo minha senha.</h2>
<p>É golpe (vishing). NINGUÉM da TI ou diretoria pede senha, por nenhum
canal. Não forneça, anote o número e reporte ao SOC.</p>

<h2>5. Recebi uma notificação de MFA que eu não solicitei.</h2>
<p>NÃO aprove. Isso indica que alguém tem sua senha (ataque de MFA fatigue).
Troque a senha imediatamente e avise o SOC no ramal 4500.</p>

<h2>6. Encontrei um pendrive no estacionamento.</h2>
<p>NÃO conecte em nenhum equipamento — é uma técnica clássica de ataque
(baiting). Entregue ao SOC sem conectar.</p>

<h2>7. Meu computador está lento e com pop-ups estranhos.</h2>
<p>Pode ser malware. Desconecte da rede, não desligue e abra chamado
urgente com o SOC (ramal 4500).</p>

<h2>8. Enviei um e-mail com dados confidenciais para a pessoa errada.</h2>
<p>Isso é um incidente de vazamento de dados. Avise o SOC imediatamente e,
se envolver dados pessoais, o SOC acionará o DPO. Não tente "resolver"
pedindo para o destinatário apagar.</p>

<h2>9. Quero instalar um programa que não está na loja corporativa.</h2>
<p>Não instale (shadow IT é proibido). Solicite homologação pelo helpdesk —
prazo médio de avaliação: 5 dias úteis.</p>
</body>
</html>"""
with open("docs/faq_incidentes.html", "w", encoding="utf-8") as f:
    f.write(faq_html)
print("✅ faq_incidentes.html")

print("\n🎉 Todos os 8 documentos de Segurança gerados na pasta docs/!")